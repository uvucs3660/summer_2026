#!/usr/bin/env python3
"""Render CS 3540 lecture deck markdown into editable .pptx with speaker notes.

Source format (see README.md): YAML frontmatter, slides separated by a line of
`---`, an optional `# Title`, body lines, and a trailing `NOTES:` block whose
text becomes the PowerPoint speaker note — that note is the recording script.

Diagrams are referenced as SVG and rasterized through rsvg-convert into
_render/; python-pptx cannot embed SVG directly.
"""
import re
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
# One diagram library, shared by the decks and the Canvas cheat sheets. Art authored
# for a lecture goes here too, so students get it on the cheat sheet page.
DIAGRAM_DIRS = [HERE.parent.parent / "cheatsheets" / "diagrams"]
CACHE = HERE / "_render"

BG = RGBColor(0x0B, 0x12, 0x20)
AMBER = RGBColor(0xFC, 0xD3, 0x4D)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
BODY = RGBColor(0xE5, 0xE7, 0xEB)
MUTE = RGBColor(0x9C, 0xA3, 0xAF)
CODE_FG = RGBColor(0x93, 0xC5, 0xFD)
CODE_BG = RGBColor(0x11, 0x18, 0x27)
GREEN = RGBColor(0x34, 0xD3, 0x99)

SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.75
BODY_TOP, BODY_BOTTOM = 1.72, 6.72

MONO = "Menlo"
SANS = "Helvetica Neue"


# ---------------------------------------------------------------- parsing
def split_slides(raw: str):
    """Split on a lone `---`, but never one inside a fenced code block.

    Decks routinely show YAML frontmatter as an example, and that frontmatter's
    own `---` delimiters would otherwise cut the slide in half.
    """
    chunks, cur, fenced = [], [], False
    for line in raw.split("\n"):
        if line.strip().startswith("```"):
            fenced = not fenced
        if not fenced and line.strip() == "---":
            chunks.append("\n".join(cur))
            cur = []
        else:
            cur.append(line)
    chunks.append("\n".join(cur))
    return chunks


def parse_deck(path: Path):
    raw = path.read_text()
    meta = {}
    if raw.startswith("---\n"):
        fm, raw = raw[4:].split("\n---\n", 1)
        for line in fm.splitlines():
            if ":" in line:
                k, v = line.split(":", 1)
                meta[k.strip()] = v.strip().strip('"')

    slides = []
    for chunk in split_slides(raw):
        if not chunk.strip():
            continue
        body_src, _, notes = chunk.partition("\nNOTES:")
        slides.append({"blocks": parse_blocks(body_src), "notes": notes.strip()})
    return meta, slides


def parse_blocks(src: str):
    """Split slide body into typed blocks: title, bullets, code, image, quote, para."""
    blocks, lines, i = [], src.split("\n"), 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if not stripped:
            i += 1
        elif stripped.startswith("```"):
            lang = stripped[3:].strip()
            i += 1
            buf = []
            while i < len(lines) and not lines[i].strip().startswith("```"):
                buf.append(lines[i])
                i += 1
            i += 1
            blocks.append({"t": "code", "lang": lang, "lines": buf})
        elif stripped.startswith("# "):
            blocks.append({"t": "title", "text": stripped[2:].strip()})
            i += 1
        elif stripped.startswith("## "):
            blocks.append({"t": "subtitle", "text": stripped[3:].strip()})
            i += 1
        elif stripped.startswith("!["):
            m = re.search(r"\((.+?)\)", stripped)
            if m:
                blocks.append({"t": "image", "src": m.group(1)})
            i += 1
        elif stripped.startswith("> "):
            blocks.append({"t": "quote", "text": stripped[2:].strip()})
            i += 1
        elif stripped.startswith("- "):
            items = []
            while i < len(lines) and (lines[i].strip().startswith("- ") or not lines[i].strip()):
                if lines[i].strip().startswith("- "):
                    depth = 1 if lines[i].startswith("  ") else 0
                    items.append((depth, lines[i].strip()[2:].strip()))
                    i += 1
                else:
                    break
            blocks.append({"t": "bullets", "items": items})
        else:
            blocks.append({"t": "para", "text": stripped})
            i += 1
    return blocks


# ------------------------------------------------------------ inline runs
def add_runs(para, text, size, color, bold=False):
    """Render **bold** and `code` as separate runs inside one paragraph."""
    for part in re.split(r"(\*\*.+?\*\*|`.+?`)", text):
        if not part:
            continue
        run = para.add_run()
        f = run.font
        f.size, f.bold = Pt(size), bold
        if part.startswith("**") and part.endswith("**"):
            run.text, f.bold, f.color.rgb = part[2:-2], True, AMBER
        elif part.startswith("`") and part.endswith("`"):
            run.text, f.name, f.color.rgb = part[1:-1], MONO, CODE_FG
            f.size = Pt(size - 2)
        else:
            run.text, f.name, f.color.rgb = part, SANS, color


# ---------------------------------------------------------------- drawing
def rasterize(src: str) -> Path:
    name = Path(src).name
    for d in DIAGRAM_DIRS:
        svg = d / name
        if svg.exists():
            break
    else:
        raise FileNotFoundError(
            f"diagram {name!r} not in " + " or ".join(str(d) for d in DIAGRAM_DIRS)
        )
    CACHE.mkdir(exist_ok=True)
    png = CACHE / (svg.stem + ".png")
    if not png.exists() or png.stat().st_mtime < svg.stat().st_mtime:
        subprocess.run(
            ["rsvg-convert", "-w", "2000", "-f", "png", "-o", str(png), str(svg)],
            check=True,
        )
    return png


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.paragraphs[0].text = ""
    return tf


def draw_title_slide(prs, meta):
    s = blank(prs)
    track = meta.get("track", "").upper()
    tf = textbox(s, MARGIN, 2.25, SLIDE_W - 2 * MARGIN, 3.0)

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"WEEK {meta.get('week', '?')}  ·  {track} TRACK"
    r.font.size, r.font.name, r.font.bold = Pt(16), SANS, True
    r.font.color.rgb = GREEN

    p = tf.add_paragraph()
    p.space_before = Pt(14)
    r = p.add_run()
    r.text = meta.get("title", "")
    r.font.size, r.font.name, r.font.bold = Pt(54), SANS, True
    r.font.color.rgb = AMBER

    if meta.get("subtitle"):
        p = tf.add_paragraph()
        p.space_before = Pt(8)
        r = p.add_run()
        r.text = meta["subtitle"]
        r.font.size, r.font.name = Pt(26), SANS
        r.font.color.rgb = BLUE

    tf2 = textbox(s, MARGIN, 6.35, SLIDE_W - 2 * MARGIN, 0.5)
    r = tf2.paragraphs[0].add_run()
    rt = meta.get("runtime")
    r.text = "CS 3540 · Game Programming · Fall 2026 · Professor Hunter" + (f"  ·  ~{rt} min" if rt else "")
    r.font.size, r.font.name = Pt(13), SANS
    r.font.color.rgb = MUTE
    return s


def draw_footer(slide, meta, n):
    tf = textbox(slide, MARGIN, 6.92, SLIDE_W - 2 * MARGIN, 0.4)
    p = tf.paragraphs[0]
    r = p.add_run()
    track = meta.get("track", "")
    track = track.upper() if len(track) <= 2 else track.title()   # "AI", not "Ai"
    r.text = f"CS 3540 · W{meta.get('week','?')} {track}"
    r.font.size, r.font.name = Pt(10), SANS
    r.font.color.rgb = MUTE
    r2 = p.add_run()
    r2.text = f"{'':>4}·{'':>4}{n}"
    r2.font.size, r2.font.name = Pt(10), SANS
    r2.font.color.rgb = MUTE
    p.alignment = PP_ALIGN.LEFT


def body_font_size(blocks):
    """Shrink type as a slide gets busier, so nothing overflows the frame."""
    weight = 0
    for b in blocks:
        if b["t"] == "bullets":
            weight += sum(1 + len(txt) // 60 for _, txt in b["items"])
        elif b["t"] == "code":
            weight += len(b["lines"])
        elif b["t"] in ("para", "quote"):
            weight += 1 + len(b["text"]) // 60
    if weight <= 5:
        return 24
    if weight <= 8:
        return 21
    if weight <= 12:
        return 18
    return 16


def draw_content(prs, meta, slide_def, n):
    s = blank(prs)
    blocks = slide_def["blocks"]

    title = next((b for b in blocks if b["t"] == "title"), None)
    rest = [b for b in blocks if b["t"] != "title"]

    if title:
        tf = textbox(s, MARGIN, 0.5, SLIDE_W - 2 * MARGIN, 1.0)
        p = tf.paragraphs[0]
        add_runs(p, title["text"], 34, AMBER, bold=True)
        for r in p.runs:
            if r.font.color.rgb != CODE_FG:
                r.font.color.rgb = AMBER
        rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(1.48), Inches(1.6), Inches(0.045))
        rule.fill.solid()
        rule.fill.fore_color.rgb = GREEN
        rule.line.fill.background()
        rule.shadow.inherit = False

    top = BODY_TOP if title else 1.0
    size = body_font_size(rest)
    images = [b for b in rest if b["t"] == "image"]
    text_blocks = [b for b in rest if b["t"] != "image"]

    # Text occupies the full width unless a diagram shares the slide.
    text_w = SLIDE_W - 2 * MARGIN
    if images and text_blocks:
        text_w = 5.1

    if text_blocks:
        tf = textbox(s, MARGIN, top, text_w, BODY_BOTTOM - top, anchor=MSO_ANCHOR.MIDDLE)
        first = True
        for b in text_blocks:
            first = draw_text_block(s, tf, b, size, first)

    if images:
        place_images(s, images, bool(text_blocks), top)

    draw_footer(s, meta, n)
    if slide_def["notes"]:
        s.notes_slide.notes_text_frame.text = slide_def["notes"]
    return s


def png_size(path: Path):
    """Width/height straight from the IHDR chunk — avoids a Pillow dependency."""
    import struct

    with path.open("rb") as fh:
        head = fh.read(24)
    if head[:8] != b"\x89PNG\r\n\x1a\n":
        raise ValueError(f"not a PNG: {path}")
    return struct.unpack(">II", head[16:24])


def place_images(slide, images, beside_text, top):
    avail_w = (SLIDE_W - MARGIN - 5.95) if beside_text else (SLIDE_W - 2 * MARGIN)
    left_edge = 5.95 if beside_text else MARGIN
    avail_h = BODY_BOTTOM - top
    if len(images) > 1:
        avail_h = (avail_h - 0.25 * (len(images) - 1)) / len(images)

    y = top
    for img in images:
        png = rasterize(img["src"])
        pw, ph = png_size(png)
        scale = min(avail_w / pw, avail_h / ph)
        w, h = pw * scale, ph * scale
        slide.shapes.add_picture(
            str(png),
            Inches(left_edge + (avail_w - w) / 2),
            Inches(y + (avail_h - h) / 2),
            Inches(w),
            Inches(h),
        )
        y += avail_h + 0.25


def draw_text_block(slide, tf, b, size, first):
    def para():
        nonlocal first
        if first:
            first = False
            return tf.paragraphs[0]
        return tf.add_paragraph()

    if b["t"] == "subtitle":
        p = para()
        p.space_before = Pt(10)
        add_runs(p, b["text"], size + 3, BLUE, bold=True)
        for r in p.runs:
            r.font.color.rgb = BLUE
    elif b["t"] == "bullets":
        for depth, text in b["items"]:
            p = para()
            p.level = depth
            p.space_before = Pt(9 if depth == 0 else 4)
            r = p.add_run()
            r.text = "▸  " if depth == 0 else "–  "
            r.font.size, r.font.name = Pt(size), SANS
            r.font.color.rgb = GREEN if depth == 0 else MUTE
            add_runs(p, text, size - (2 if depth else 0), BODY)
    elif b["t"] == "para":
        p = para()
        p.space_before = Pt(10)
        add_runs(p, b["text"], size, BODY)
    elif b["t"] == "quote":
        p = para()
        p.space_before = Pt(12)
        add_runs(p, b["text"], size, AMBER)
        for r in p.runs:
            r.font.italic = True
    elif b["t"] == "code":
        for ln in b["lines"]:
            p = para()
            p.space_before = Pt(2)
            r = p.add_run()
            r.text = ln if ln.strip() else " "
            r.font.size, r.font.name = Pt(max(13, size - 4)), MONO
            r.font.color.rgb = CODE_FG
    return first


def build(src: Path, out: Path):
    meta, slides = parse_deck(src)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SLIDE_W), Inches(SLIDE_H)

    first, rest = slides[0], slides[1:]
    s = draw_title_slide(prs, meta)
    if first["notes"]:
        s.notes_slide.notes_text_frame.text = first["notes"]
    for n, sd in enumerate(rest, start=2):
        draw_content(prs, meta, sd, n)

    prs.save(str(out))
    words = sum(len(sd["notes"].split()) for sd in slides)
    return len(slides), words


def main(argv):
    targets = [Path(a) for a in argv[1:]] or sorted(HERE.glob("w*-*.md"))
    if not targets:
        print("no deck sources found", file=sys.stderr)
        return 1
    for src in targets:
        out = src.with_suffix(".pptx")
        n, words = build(src, out)
        # ~140 wpm is a normal recorded-lecture pace.
        print(f"  {out.name}: {n} slides, {words} script words (~{words // 140} min)")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
